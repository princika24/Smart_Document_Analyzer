import pdfplumber
import docx
from pptx import Presentation
import re

def extract_text(file_path: str) -> str:
    ext = file_path.split(".")[-1].lower()
    if ext == "pdf":
        text = _extract_pdf(file_path)
    elif ext == "docx":
        text = _extract_docx(file_path)
    elif ext == "pptx":
        text = _extract_pptx(file_path)
    elif ext == "txt":
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
    else:
        raise ValueError("Unsupported file type")

    return _clean_text(text)


def _extract_pdf(path):
    text = ""
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
    return text


def _extract_docx(path):
    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx(path):
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text)
    return "\n".join(texts)


def _clean_text(text: str) -> str:
    """Clean noisy text extracted from documents to prevent hallucination."""
    if not text:
        return ""

    # Remove URLs, emails, weird site-like text
    text = re.sub(r"http\S+|www\S+|@\S+", " ", text)

    # Remove typical spammy or navigation words
    text = re.sub(
        r"\b(share|photo|video|page|click|email|home|comment|subscribe|story|article|report)\b",
        " ",
        text,
        flags=re.IGNORECASE,
    )

    # Remove duplicate spaces, fix formatting
    text = re.sub(r"\s+", " ", text).strip()
    return text
